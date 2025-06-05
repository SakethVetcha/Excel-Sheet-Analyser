import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import gc
from io import BytesIO
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
import numpy as np

@st.cache_data
def load_excel_sheet(_file, sheet_name):
    """Cache the loading of individual sheets to prevent reloading"""
    try:
        # Determine the engine based on the file object's name if available, otherwise use openpyxl
        if hasattr(_file, 'name'):
            engine = 'openpyxl' if _file.name.endswith('.xlsx') else 'pyxlsb'
        else:
            # Default to openpyxl if we can't determine the file type
            engine = 'openpyxl'
            
        # Read Excel file with optimized memory usage
        df = pd.read_excel(
            _file, 
            sheet_name=sheet_name,
            engine=engine
        )
        
        # Optimize memory usage
        for col in df.columns:
            if df[col].dtype == 'object':
                # Convert object types to categories if they have few unique values
                if df[col].nunique() / len(df) < 0.5:  # If less than 50% unique values
                    df[col] = df[col].astype('category')
            elif df[col].dtype == 'float64':
                # Downcast float64 to float32 if possible
                df[col] = pd.to_numeric(df[col], downcast='float')
            elif df[col].dtype == 'int64':
                # Downcast int64 to smallest integer type possible
                df[col] = pd.to_numeric(df[col], downcast='integer')
        
        return df
    except Exception as e:
        st.error(f"Error loading sheet {sheet_name}: {str(e)}")
        return None

class FlexibleDataAnalysis:
    def __init__(self, df):
        self.df = df
        
        self.has_type = 'Type' in df.columns
        self.has_price = 'Price' in df.columns
        self.has_quantity = 'Quantity' in df.columns
        
        if not self.has_price:
            st.error("Error: 'Price' column is missing in the data. Please ensure you've selected the correct Price column.")
            return
            
        # Convert all string columns to lowercase
        for col in self.df.columns:
            if self.df[col].dtype == 'object':
                self.df[col] = self.df[col].str.lower()
            
        # Convert Price and Quantity columns with memory optimization
        if self.has_price:
            self.df['Price'] = pd.to_numeric(self.df['Price'], errors='coerce', downcast='float')
        
        if self.has_quantity:
            self.df['Quantity'] = pd.to_numeric(self.df['Quantity'], errors='coerce', downcast='integer')
        else:
            # If no Quantity column, default to 1
            self.df['Quantity'] = 1
            self.has_quantity = True
        
        # Calculate total price with optimized types
        try:
            self.df['Total_Price'] = self.df['Price'] * self.df['Quantity']
            self.df['Total_Price'] = pd.to_numeric(self.df['Total_Price'], downcast='float')
        except Exception as e:
            st.error(f"Error calculating Total Price: {str(e)}")
            self.df['Total_Price'] = 0

    def create_pie_chart(self, column, title):
        if column not in self.df.columns:
            return None
        
        plt.figure(figsize=(15, 10))  # Increased figure size
        
        # Convert the column values to lowercase and then get value counts
        data = self.df[column].str.lower().value_counts()
        
        # If too many categories, group small ones into "Others"
        if len(data) > 15:
            threshold = data.sum() * 0.01  # 1% threshold
            other_mask = data < threshold
            if other_mask.any():
                other_sum = data[other_mask].sum()
                data = data[~other_mask]
                data['others'] = other_sum
        
        # Sort values in descending order
        data = data.sort_values(ascending=False)
        
        # Calculate percentages
        total = data.sum()
        percentages = [(val/total)*100 for val in data.values]
        
        # Generate colors using a color map
        colors = plt.cm.Set3(np.linspace(0, 1, len(data)))
        
        # Create pie chart
        plt.pie(
            data.values,
            labels=None,
            colors=colors,
            startangle=90,
            radius=1.2,
            counterclock=False,
            wedgeprops={'linewidth': 2, 'edgecolor': 'white'}
        )
        
        # Create custom legend with color boxes
        legend_elements = []
        for i, (label, percentage) in enumerate(zip(data.index, percentages)):
            legend_elements.append(plt.Rectangle((0, 0), 1, 1, fc=colors[i], label=f'{label} ({percentage:.1f}%)'))
        
        # Add legend with color boxes
        plt.legend(
            handles=legend_elements,
            loc='center left',
            bbox_to_anchor=(1.1, 0.5),
            fontsize=12,
            title='Categories',
            title_fontsize=14,
            frameon=True,
            facecolor='white',
            edgecolor='gray',
            framealpha=0.9,
            borderpad=1
        )
        
        # Add title with padding
        plt.title(title, pad=20, size=18, weight='bold')
        
        # Equal aspect ratio ensures circular pie
        plt.axis('equal')
        
        # Add more padding around the entire figure
        plt.tight_layout(pad=3.0)
        
        return plt.gcf()

    def create_all_pie_charts(self):
        charts = {}
        
        with st.spinner('Creating pie charts...'):
            if 'Source' in self.df.columns:
                charts['Source of Scan'] = self.create_pie_chart('Source', 'Distribution by Source of Scan')
            if 'Status' in self.df.columns:
                charts['Old/New'] = self.create_pie_chart('Status', 'Distribution by Old/New Status')
            if 'Transaction_Status' in self.df.columns:
                charts['Transaction Status'] = self.create_pie_chart('Transaction_Status', 'Distribution by Transaction Status')
            if 'Payment_Mode' in self.df.columns:
                charts['Payment Mode'] = self.create_pie_chart('Payment_Mode', 'Distribution by Payment Mode')
            if 'Product_Name' in self.df.columns:
                charts['Product Name'] = self.create_pie_chart('Product_Name', 'Distribution by Product Name')
        
        return charts

    def basic_statistics(self):
        if self.df is None:
            return None
        stats = {}
        
        try:
            total_revenue = self.df['Total_Price'].sum()
            avg_unit_price = self.df['Price'].mean()
            max_unit_price = self.df['Price'].max()
            min_unit_price = self.df['Price'].min()
            total_quantity = self.df['Quantity'].sum()
            avg_transaction_value = total_revenue / len(self.df)
            
            stats["Total Revenue"] = f"${total_revenue:,.2f}" if pd.notnull(total_revenue) else "N/A"
            stats["Average Unit Price"] = f"${avg_unit_price:,.2f}" if pd.notnull(avg_unit_price) else "N/A"
            stats["Highest Unit Price"] = f"${max_unit_price:,.2f}" if pd.notnull(max_unit_price) else "N/A"
            stats["Lowest Unit Price"] = f"${min_unit_price:,.2f}" if pd.notnull(min_unit_price) else "N/A"
            stats["Total Quantity Sold"] = f"{total_quantity:,.0f}" if pd.notnull(total_quantity) else "N/A"
            stats["Average Transaction Value"] = f"${avg_transaction_value:,.2f}" if pd.notnull(avg_transaction_value) else "N/A"
            stats["Total Unique Items"] = str(len(self.df['Item'].unique()))
            stats["Total Transactions"] = str(len(self.df))
            
            if self.has_type:
                stats["Number of Categories"] = str(len(self.df['Type'].unique()))
        except Exception as e:
            st.error(f"Error calculating statistics: {str(e)}")
            return pd.Series({"Error": "Could not calculate statistics"})
        
        return pd.Series(stats)

    def generate_presentation(self, title="Data Analysis Report"):
        if self.df is None:
            return None
            
        try:
            prs = Presentation()
            
            # Set slide dimensions to 16:9 aspect ratio
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            # Add title safely
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = title
            else:
                left = Inches(1)
                top = Inches(1)
                width = Inches(11.333)  # Adjusted for 16:9
                height = Inches(1.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = title

            # Add data and time to the first slide 
            try:
                subtitle = slide.placeholders[1]
                subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            except:
                left = Inches(1)
                top = Inches(2.5)
                width = Inches(11.333)  # Adjusted for 16:9
                height = Inches(1)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

            # Basic Statistics slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = 'Basic Statistics'
            else:
                left = Inches(1)
                top = Inches(0.5)
                width = Inches(11.333)  # Adjusted for 16:9
                height = Inches(1)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = 'Basic Statistics'

            # Add statistics content
            stats = self.basic_statistics()
            try:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
            except:
                left = Inches(1)
                top = Inches(2)
                width = Inches(11.333)  # Adjusted for 16:9
                height = Inches(4)
                txBox = slide.shapes.add_textbox(left, top, width, height)  
                tf = txBox.text_frame

            # Create two columns for statistics
            stats_items = list(stats.items())
            mid_point = len(stats_items) // 2
            
            # Left column
            left_col = tf.add_paragraph()
            for stat_name, stat_value in stats_items[:mid_point]:
                p = tf.add_paragraph()
                p.text = f"{stat_name}: {stat_value}"
                p.level = 0
            
            # Right column (if needed)
            if mid_point < len(stats_items):
                right_col_box = slide.shapes.add_textbox(
                    left=Inches(7),
                    top=Inches(2),
                    width=Inches(5.333),
                    height=Inches(4)
                )
                right_tf = right_col_box.text_frame
                for stat_name, stat_value in stats_items[mid_point:]:
                    p = right_tf.add_paragraph()
                    p.text = f"{stat_name}: {stat_value}"

            # Pie Charts slides
            charts = self.create_all_pie_charts()
            if charts:
                for chart_title, fig in charts.items():
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                    
                    # Add title
                    left = Inches(1)
                    top = Inches(0.5)
                    width = Inches(11.333)  # Adjusted for 16:9
                    height = Inches(1)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = chart_title
                    
                    # Adjust figure size for 16:9 aspect ratio and add padding
                    plt.figure(figsize=(15, 10))
                    plt.subplots_adjust(left=0.15, right=0.85, top=0.85, bottom=0.15)
                    
                    # Save the matplotlib figure to a BytesIO object with high quality
                    img_stream = BytesIO()
                    fig.savefig(
                        img_stream, 
                        format='png', 
                        bbox_inches='tight', 
                        dpi=300, 
                        pad_inches=1.0,  # Increased padding
                        facecolor='white'
                    )
                    img_stream.seek(0)
                    
                    # Calculate dimensions to maintain aspect ratio while leaving space for labels
                    img_width = Inches(11)  # Wider to accommodate larger text
                    img_height = Inches(6.5)  # Taller to accommodate larger text
                    
                    # Center the image on the slide
                    left = (prs.slide_width - img_width) / 2
                    top = Inches(0.75)  # Move up slightly to fit larger chart
                    pic = slide.shapes.add_picture(img_stream, left, top, width=img_width, height=img_height)
                    plt.close(fig)

            # Save presentation to BytesIO
            output = BytesIO()
            prs.save(output)
            output.seek(0)
            return output

        except Exception as e:
            st.error(f"Error generating presentation: {str(e)}")
            return None

def main():
    st.set_page_config(page_title="Flexible Data Analyzer", layout="wide")
    st.title("📊 Data Analysis Dashboard")

    uploaded_file = st.file_uploader("Upload your Excel file", type=["xlsx", "xlsb"])

    if not uploaded_file:
        st.info("Please upload an Excel file to begin analysis.")
        return

    try:
        # Get file size and show info
        file_size = uploaded_file.size / (1024 * 1024)  # Convert to MB
        st.info(f"File size: {file_size:.2f} MB")
        
        if file_size > 100:  # If file is larger than 100MB
            st.warning("Large file detected. Processing may take a few moments...")
        
        # Create tabs for overall analysis and sheet-wise analysis
        tab1, tab2 = st.tabs(["Overall Analysis", "Sheet-wise Analysis"])
        
        # Dictionary to store all dataframes
        all_dfs = {}
        combined_df = pd.DataFrame()
        
        # Read all sheets with progress indication
        with st.spinner('Loading Excel sheets...'):
            xls = pd.ExcelFile(uploaded_file, engine='openpyxl' if uploaded_file.name.endswith('.xlsx') else 'pyxlsb')
            progress_bar = st.progress(0)
            
            for idx, sheet in enumerate(xls.sheet_names):
                df = load_excel_sheet(xls, sheet)
                if df is not None:
                    all_dfs[sheet] = df
                    combined_df = pd.concat([combined_df, df], ignore_index=True)
                progress_bar.progress((idx + 1) / len(xls.sheet_names))
            
            progress_bar.empty()
        
        # Free up memory
        gc.collect()

        # Initialize column mapping dictionary at a higher scope
        col_map = {}
        
        def process_tab1():
            nonlocal col_map
            st.subheader("Overall Analysis")
            st.info("Map your columns to the required fields:")
            columns = combined_df.columns.tolist()
            
            # Required columns
            st.warning("Please ensure you select the correct Price column. This is required for the analysis.")
            price_col = st.selectbox("Select the Price column (required)", [""] + columns, key="price_col")
            item_col = st.selectbox("Select the Item column (required)", [""] + columns, key="item_col")
            
            if not price_col:
                st.error("Please select a Price column to continue with the analysis.")
                return
            
            if not item_col:
                st.error("Please select an Item column to continue with the analysis.")
                return
            
            # Optional columns
            st.info("The following columns are optional:")
            source_col = st.selectbox("Select the Source of Scan column (optional)", ["None"] + columns, key="source_col")
            status_col = st.selectbox("Select the Old/New Status column (optional)", ["None"] + columns, key="status_col")
            trans_status_col = st.selectbox("Select the Transaction Status column (optional)", ["None"] + columns, key="trans_status_col")
            payment_col = st.selectbox("Select the Payment Mode column (optional)", ["None"] + columns, key="payment_col")
            product_name_col = st.selectbox("Select the Product Name column (optional)", ["None"] + columns, key="product_name_col")
            quantity_col = st.selectbox("Select the Quantity column (optional)", ["None"] + columns, key="quantity_col")

            if price_col and item_col:
                try:
                    # Update the column mapping
                    col_map.clear()
                    col_map.update({
                        item_col: 'Item',
                        price_col: 'Price'
                    })
                    
                    if source_col != "None":
                        col_map[source_col] = 'Source'
                    if status_col != "None":
                        col_map[status_col] = 'Status'
                    if trans_status_col != "None":
                        col_map[trans_status_col] = 'Transaction_Status'
                    if payment_col != "None":
                        col_map[payment_col] = 'Payment_Mode'
                    if product_name_col != "None":
                        col_map[product_name_col] = 'Product_Name'
                    if quantity_col != "None":
                        col_map[quantity_col] = 'Quantity'
                    
                    # Verify the Price column contains numeric data
                    try:
                        test_price = pd.to_numeric(combined_df[price_col], errors='coerce')
                        if test_price.isna().all():
                            st.error(f"The selected Price column '{price_col}' does not contain any valid numeric values. Please select a different column.")
                            return
                    except Exception as e:
                        st.error(f"Error validating Price column: {str(e)}")
                        return
                    
                    df_renamed = combined_df.rename(columns=col_map)
                    
                    # If quantity is not provided, default to 1 for each row
                    if 'Quantity' not in df_renamed.columns:
                        df_renamed['Quantity'] = 1
                    
                    analyzer = FlexibleDataAnalysis(df_renamed)
                    
                    # Basic statistics
                    st.subheader("Basic Statistics")
                    stats = analyzer.basic_statistics()
                    if stats is not None:
                        st.table(stats)
                    
                    # Generate and offer PowerPoint download
                    st.subheader("Download PowerPoint Presentation")
                    st.info("The PowerPoint presentation includes basic statistics and distribution charts.")
                    pptx = analyzer.generate_presentation()
                    if pptx is not None:
                        st.download_button(
                            label="Download Analysis PowerPoint",
                            data=pptx,
                            file_name="analysis_presentation.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                except Exception as e:
                    st.error(f"Error processing data: {str(e)}")
        
        def process_tab2():
            if not col_map:
                st.warning("Please select columns in the Overall Analysis tab first.")
                return
                
            st.subheader("Sheet-wise Analysis")
            selected_sheet = st.selectbox("Select sheet to view detailed analysis", xls.sheet_names)
            
            if selected_sheet:
                df = all_dfs[selected_sheet]
                df_renamed = df.rename(columns=col_map)
                
                if 'Quantity' not in df_renamed.columns:
                    df_renamed['Quantity'] = 1
                
                sheet_analyzer = FlexibleDataAnalysis(df_renamed)
                
                # Basic statistics for the sheet
                st.subheader("Basic Statistics")
                sheet_stats = sheet_analyzer.basic_statistics()
                if sheet_stats is not None:
                    st.table(sheet_stats)
                
                # Generate and offer PowerPoint download for individual sheet
                st.subheader("Download PowerPoint Presentation")
                st.info("The PowerPoint presentation includes basic statistics and distribution charts.")
                sheet_pptx = sheet_analyzer.generate_presentation(f"Data Analysis Report - {selected_sheet}")
                if sheet_pptx is not None:
                    st.download_button(
                        label=f"Download Analysis PowerPoint - {selected_sheet}",
                        data=sheet_pptx,
                        file_name=f"analysis_presentation_{selected_sheet}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
        
        with tab1:
            process_tab1()
        
        with tab2:
            process_tab2()
            
    except Exception as e:
        st.error(f"Error processing file: {str(e)}")

if __name__ == "__main__":
    main()

